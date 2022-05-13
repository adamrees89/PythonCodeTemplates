## Credit to Haider Imtiaz

# Automate PowerPoint
# pip install python-pptx
from pptx import Presentation
# Read Slides
PP = Presentation('file.pptx')
for slide in PP.slides:
    for shape in slide.shapes:
        for paragraph in shape.text_frame.paragraphs:
            for data in paragraph.runs:
                print(data.text)
# Write in PPT
PP = Presentation()
title_slide_layout = PP.slide_layouts[0]
slide = PP.slides.add_slide(title_slide_layout)
title = slide.shapes.title
title.text = "Medium Article"
PP.save('file.pptx')